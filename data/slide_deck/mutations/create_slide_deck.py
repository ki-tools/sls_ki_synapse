# Copyright 2018-present, Bill & Melinda Gates Foundation
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import graphene
import os
import requests
import tempfile
import boto3
import datetime
import uuid
from core import ParamStore
from core.log import logger
from ..types import SlideDeck
from pptx import Presentation


class CreateSlideDeck(graphene.Mutation):
    """
    Mutation for creating a SlideDeck.
    """
    slide_deck = graphene.Field(lambda: SlideDeck)

    class Arguments:
        title = graphene.String(required=True)
        presenter = graphene.String(required=True)
        sprint_id = graphene.String(required=True)
        participants = graphene.List(graphene.String, required=True)
        end_date = graphene.String(required=True)
        sprint_questions = graphene.List(graphene.String, required=True)
        background = graphene.String(required=True)
        problem_statement = graphene.String(required=True)
        motivation = graphene.String(required=True)
        deliverables = graphene.List(graphene.String, required=True)
        key_findings = graphene.List(graphene.String, required=True)
        next_steps = graphene.List(graphene.String, required=True)
        value = graphene.String(required=True)
        template_url = graphene.String(
            description='URL to a presentation to use as a template.')

    def mutate(self,
               info,
               title,
               presenter,
               sprint_id,
               participants,
               end_date,
               sprint_questions,
               background,
               problem_statement,
               motivation,
               deliverables,
               key_findings,
               next_steps,
               value,
               **kwargs):
        template_url = kwargs.get('template_url', None)

        presentation = None

        if template_url and template_url != '':
            logger.debug('Fetching template file: {0}'.format(template_url))
            r = requests.get(template_url)
            logger.debug('Done fetching template file.')

            if r.status_code == 200:
                fd, tmp_filename = tempfile.mkstemp(suffix='.pptx')

                try:
                    with os.fdopen(fd, 'wb') as tmp:
                        tmp.write(r.content)

                    presentation = Presentation(tmp_filename)
                finally:
                    if os.path.isfile(tmp_filename):
                        os.remove(tmp_filename)
            else:
                raise Exception(
                    'Could not load template_url: {0}'.format(template_url))
        else:
            presentation = Presentation('assets/template_ki_empty.pptx')

        SLD_TITLE = 'Title Slide - Text Only'
        SLD_HEAD_COPY = 'Full Width Head'
        SLD_HEAD_BULLETS = 'Full Width Head + Bullets'
        # SLD_HEAD_SUBHEAD_COPY = 3
        # SLD_HEAD_ONLY = 7
        SLD_INSTRUCTIONS = 'INSTRUCTIONS'

        title_layout = presentation.slide_layouts.get_by_name(SLD_TITLE)
        plain_layout = presentation.slide_layouts.get_by_name(SLD_HEAD_COPY)
        bullet_layout = presentation.slide_layouts.get_by_name(SLD_HEAD_BULLETS)

        if title_layout is None or plain_layout is None or bullet_layout is None:
            raise Exception('Slide deck provided is not using the correct template')

        # deliverables

        slide = presentation.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = 'Deliverables'
        tf = body_shape.text_frame
        for item in deliverables:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # questions

        slide = presentation.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = 'Sprint Questions'
        tf = body_shape.text_frame
        for item in sprint_questions:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # problem statement

        slide = presentation.slides.add_slide(plain_layout)
        title_shape = CreateSlideDeck.get_placeholder(slide, 'Title 2')
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = "Problem Statement"
        body_shape.text = problem_statement
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # motivation

        slide = presentation.slides.add_slide(plain_layout)
        title_shape = CreateSlideDeck.get_placeholder(slide, 'Title 2')
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = "Motivation"
        body_shape.text = motivation
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # background

        slide = presentation.slides.add_slide(plain_layout)
        title_shape = CreateSlideDeck.get_placeholder(slide, 'Title 2')
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = "Background"
        body_shape.text = background
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # title slide

        slide = presentation.slides.add_slide(title_layout)
        title_shape = CreateSlideDeck.get_placeholder(slide, 'Title 1')
        body_shape1 = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 2')
        body_shape2 = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 3')
        title_shape.text = 'Rally {0}: {1}'.format(sprint_id, title)
        body_shape1.text = 'Completed {0}'.format(end_date)
        body_shape2.text = 'Presented by {0} on behalf of rally participants {1}'.format(
            presenter, ', '.join(participants))
        CreateSlideDeck.add_notes(slide)

        CreateSlideDeck.move_to_front(presentation)

        # data/methods/results are already part of the deck

        # key findings

        slide = presentation.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = 'Key Findings'
        tf = body_shape.text_frame
        for item in key_findings:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        CreateSlideDeck.add_notes(slide)

        # value

        slide = presentation.slides.add_slide(plain_layout)
        title_shape = CreateSlideDeck.get_placeholder(slide, 'Title 2')
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = "Value"
        body_shape.text = value
        CreateSlideDeck.add_notes(slide)

        # next steps

        slide = presentation.slides.add_slide(bullet_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = CreateSlideDeck.get_placeholder(slide, 'Text Placeholder 1')
        title_shape.text = 'Next Steps'
        tf = body_shape.text_frame
        items = next_steps
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
        CreateSlideDeck.add_notes(slide)

        # remove INSTRUCTIONS slide before saving

        slides = list(presentation.slides)
        slides2 = list(presentation.slides._sldIdLst)
        rm_idx = next((i for i in range(len(slides))
                       if slides[i].slide_layout.name == SLD_INSTRUCTIONS), None)
        if rm_idx != None:
            presentation.slides._sldIdLst.remove(slides2[rm_idx])

        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S%f")
        ppt_file_name = 'Rally_{0}_report-out-{1}.pptx'.format(sprint_id, timestamp)
        ppt_file_path = os.path.join(tempfile.gettempdir(), ppt_file_name)

        presentation.save(ppt_file_path)

        # Store on S3
        try:
            logger.debug('Uploading SlideDeck to S3.')

            s3 = boto3.resource('s3')

            s3.meta.client.upload_file(ppt_file_path, ParamStore.SLIDE_DECKS_BUCKET_NAME(), ppt_file_name)

            logger.debug('Finished uploading SlideDeck to S3.')

            # Generate a presigned URL that expires.
            presigned_url = s3.meta.client.generate_presigned_url(
                'get_object',
                Params={'Bucket': ParamStore.SLIDE_DECKS_BUCKET_NAME(), 'Key': ppt_file_name},
                ExpiresIn=ParamStore.SLIDE_DECKS_URL_EXPIRES_IN_SECONDS()
            )
        finally:
            if os.path.isfile(ppt_file_path):
                os.remove(ppt_file_path)

        new_slide_deck = SlideDeck(url=presigned_url)

        return CreateSlideDeck(slide_deck=new_slide_deck)

    @staticmethod
    def move_to_front(presentation):
        last_slide = len(presentation.slides) - 1
        slides = list(presentation.slides._sldIdLst)
        presentation.slides._sldIdLst.remove(slides[last_slide])
        presentation.slides._sldIdLst.insert(0, slides[last_slide])

    @staticmethod
    def get_placeholder(slide, name):
        for shape in slide.placeholders:
            if shape.name == name:
                return shape
        raise Exception('Invalid slide deck template. Cannot find shape with name: {0}'.format(name))

    @staticmethod
    def add_notes(slide):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = 'DO NOT EDIT THIS SLIDE!!! INSTEAD, EDIT THE RALLY METADATA ON http://hub.kiglobalhealth.org.'
